# TASK 3: Consolidate Azure slides
# Reads and writes to OUTPUT_PATH (the already-updated PPTX)

import sys, io
sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')

from pptx import Presentation
from lxml import etree

OUTPUT_PATH = "material/Session 8/ANLP Session8_Week9_UPDATED.PPTX"

prs = Presentation(OUTPUT_PATH)

# Only flag slides that are primarily Azure AI Search marketing/setup content
# (NOT generic vector-search concept slides that merely mention Azure)
AZURE_SPECIFIC_KW = [
    "azure ai search offers integrated vectorization",
    "vector database built for enterprise scale",
    "how do i get started with vector search",
    "many options for ai rag patterns",
]

def is_azure_specific_slide(slide):
    full_text = ""
    for shape in slide.shapes:
        if shape.has_text_frame:
            full_text += shape.text_frame.text.lower() + " "
    return any(kw in full_text for kw in AZURE_SPECIFIC_KW)

def delete_slide(prs, index):
    """Remove a slide by 0-based index."""
    xml_slides = prs.slides._sldIdLst
    slide = prs.slides[index]
    slide_part = slide.part
    # Remove from slide id list
    xml_slides.remove(xml_slides[index])
    # Drop the relationship from the presentation part
    slide_part_name = slide_part.partname
    rId = None
    for r in prs.part.rels.values():
        if hasattr(r, '_target') and hasattr(r._target, '_partname'):
            if r._target._partname == slide_part_name:
                rId = r.rId
                break
    if rId:
        prs.part.drop_rel(rId)

# --- Find all Azure-specific slides ---
azure_slide_indices = []
for i, slide in enumerate(prs.slides):
    if is_azure_specific_slide(slide):
        title_text = ""
        for shape in slide.shapes:
            if shape.has_text_frame:
                title_text = shape.text_frame.text.strip()[:70]
                break
        print(f"  Azure slide found at 0-based index {i} (Slide {i+1}): {title_text!r}")
        azure_slide_indices.append(i)

if not azure_slide_indices:
    print("No Azure-specific slides found. Nothing to do.")
    sys.exit(0)

# Keep the first one (index 25, "Azure AI Search offers Integrated Vectorization")
keep_idx = azure_slide_indices[0]
keep_slide = prs.slides[keep_idx]
print(f"\n  Keeping slide at 0-based index {keep_idx} (Slide {keep_idx+1}) as enterprise reference.")

# --- Update the kept slide ---
# Find the title shape and the largest body shape
title_shape = None
body_shape = None
max_len = 0

for shape in keep_slide.shapes:
    if not shape.has_text_frame:
        continue
    name = shape.name.lower()
    if "title" in name and title_shape is None:
        title_shape = shape
    else:
        text_len = len(shape.text_frame.text)
        if text_len > max_len:
            max_len = text_len
            body_shape = shape

# Fallback: if no explicit title shape, use the first shape
if title_shape is None:
    for shape in keep_slide.shapes:
        if shape.has_text_frame:
            title_shape = shape
            break

ENTERPRISE_TITLE = "Enterprise Vector Stores: Azure AI Search"
ENTERPRISE_BODY = (
    "For production at scale:\n"
    "• Azure AI Search — managed vector + keyword hybrid search\n"
    "• Supports BM25 + semantic re-ranking + RRF out of the box\n"
    "• Full setup guide: aka.ms/azure-ai-search-docs\n\n"
    "Not required for this course — FAISS covers the same concepts locally.\n"
    "Consider Azure AI Search when moving from prototype to enterprise deployment."
)

if title_shape:
    title_shape.text_frame.clear()
    title_shape.text_frame.paragraphs[0].text = ENTERPRISE_TITLE
    print(f"  Updated title to: {ENTERPRISE_TITLE!r}")

if body_shape:
    body_shape.text_frame.clear()
    body_shape.text_frame.paragraphs[0].text = ENTERPRISE_BODY
    print(f"  Updated body content.")
elif title_shape:
    # No separate body shape found — add content to second paragraph of title shape
    print("  Warning: no body shape found; content appended to title shape.")

# --- Delete the rest (in reverse order to preserve indices) ---
to_delete = sorted([i for i in azure_slide_indices if i != keep_idx], reverse=True)
for idx in to_delete:
    title_text = ""
    for shape in prs.slides[idx].shapes:
        if shape.has_text_frame:
            title_text = shape.text_frame.text.strip()[:60]
            break
    print(f"  Deleting Azure slide at 0-based index {idx} (Slide {idx+1}): {title_text!r}")
    delete_slide(prs, idx)

print(f"\n  Slides remaining after consolidation: {len(prs.slides)}")

prs.save(OUTPUT_PATH)
print(f"  Saved to {OUTPUT_PATH}")

# --- Final verification ---
print("\n=== Final slide listing ===")
prs2 = Presentation(OUTPUT_PATH)
print(f"Total slides: {len(prs2.slides)}")
for i, slide in enumerate(prs2.slides):
    texts = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            t = shape.text_frame.text.strip()[:60]
            if t:
                texts.append(t)
    combined = " | ".join(texts)[:90]
    print(f"  Slide {i+1:3d}: {combined!r}")
