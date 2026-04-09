# material/Session 8/update_pptx.py
from pptx import Presentation

INPUT_PATH = "material/Session 8/ANLP Session8_Week9.PPTX"
OUTPUT_PATH = "material/Session 8/ANLP Session8_Week9_UPDATED.PPTX"

prs = Presentation(INPUT_PATH)


def delete_slide(prs, index):
    """Delete slide at 0-based index."""
    xml_slides = prs.slides._sldIdLst
    slide_part = prs.slides[index].part
    xml_slides.remove(xml_slides[index])
    slide_part_name = slide_part.partname
    rId = None
    for r in prs.part.rels.values():
        if hasattr(r, '_target') and hasattr(r._target, '_partname'):
            if r._target._partname == slide_part_name:
                rId = r.rId
                break
    if rId:
        prs.part.drop_rel(rId)


# 1-based slide numbers to delete (from audit):
# - Slides 22-40: Prompt Engineering section (19 slides)
# - Slides 60-66: Azure setup tutorial (7 slides)
# - Slide 14: MCP slide
slides_to_delete_1based = list(range(22, 41)) + list(range(60, 67)) + [14]

# Convert to 0-based and sort DESCENDING (delete from end first to preserve indices)
slides_to_delete_0based = sorted(set(s - 1 for s in slides_to_delete_1based), reverse=True)

print(f"Slides before deletion: {len(prs.slides)}")
for idx in slides_to_delete_0based:
    if idx < len(prs.slides):
        print(f"  Deleting 0-based index {idx}")
        delete_slide(prs, idx)

print(f"Slides after deletion: {len(prs.slides)}")
prs.save(OUTPUT_PATH)
print(f"Saved to {OUTPUT_PATH}")
