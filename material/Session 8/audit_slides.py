# material/Session 8/audit_slides.py
from pptx import Presentation

pptx_path = "material/Session 8/ANLP Session8_Week9.PPTX"
prs = Presentation(pptx_path)

for i, slide in enumerate(prs.slides):
    title = ""
    body = ""
    for shape in slide.shapes:
        if shape.has_text_frame:
            text = shape.text_frame.text.strip()
            if not title and ("title" in shape.name.lower()):
                title = text[:80]
            elif not body:
                body = text[:80]
    if not title:
        for shape in slide.shapes:
            if shape.has_text_frame:
                title = shape.text_frame.text.strip()[:80]
                break
    print(f"Slide {i+1:3d}: {title!r}")
    if body and body != title:
        print(f"         {body!r}")
